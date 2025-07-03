"""Core presentation builder functionality."""

from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches

from .templates import PresentationTemplate, SlideTemplate, SlideType


class PresentationBuilder:
    """Main class for building presentations programmatically."""

    def __init__(self, template_path: Path | None = None):
        """Initialize the presentation builder.

        Args:
            template_path: Path to PowerPoint template file (.pptx)
        """
        if template_path and template_path.exists():
            self.presentation = Presentation(str(template_path))
        else:
            self.presentation = Presentation()

        self.slides: list[Slide] = []

    def add_title_slide(self, title: str, subtitle: str | None = None) -> Slide:
        """Add a title slide to the presentation.

        Args:
            title: Main title text
            subtitle: Optional subtitle text

        Returns:
            The created slide
        """
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle

        self.slides.append(slide)
        return slide  # type: ignore[no-any-return]

    def add_content_slide(self, title: str, content: list[str]) -> Slide:
        """Add a content slide with bullet points.

        Args:
            title: Slide title
            content: List of bullet point texts

        Returns:
            The created slide
        """
        slide_layout = self.presentation.slide_layouts[1]  # Content slide layout
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()

        for i, item in enumerate(content):
            paragraph = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            paragraph.text = item
            paragraph.level = 0

        self.slides.append(slide)
        return slide  # type: ignore[no-any-return]

    def add_image_slide(self, title: str, image_path: Path,
                       caption: str | None = None) -> Slide:
        """Add a slide with an image.

        Args:
            title: Slide title
            image_path: Path to image file
            caption: Optional image caption

        Returns:
            The created slide
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_shape.text = title

        # Add image
        if image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(1), Inches(1.5),
                width=Inches(8)
            )

        # Add caption if provided
        if caption:
            caption_shape = slide.shapes.add_textbox(
                Inches(1), Inches(6.5), Inches(8), Inches(1)
            )
            caption_shape.text = caption

        self.slides.append(slide)
        return slide  # type: ignore[no-any-return]

    def add_section_slide(self, title: str, description: str | None = None) -> Slide:
        """Add a section divider slide.

        Args:
            title: Section title
            description: Optional section description

        Returns:
            The created slide
        """
        slide_layout = self.presentation.slide_layouts[2]  # Section header layout
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if description and len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            content_shape.text = description

        self.slides.append(slide)
        return slide  # type: ignore[no-any-return]

    def build_from_template(self, template: PresentationTemplate) -> None:
        """Build presentation from a template.

        Args:
            template: Presentation template with slide definitions
        """
        for slide_template in template.slides:
            self._add_slide_from_template(slide_template)

    def _add_slide_from_template(self, slide_template: SlideTemplate) -> Slide:
        """Add a slide from a slide template.

        Args:
            slide_template: Template defining the slide

        Returns:
            The created slide
        """
        slide_type = slide_template.slide_type
        params = slide_template.params

        if slide_type == SlideType.TITLE:
            return self.add_title_slide(
                params.get('title', ''),
                params.get('subtitle', '')
            )
        elif slide_type == SlideType.CONTENT:
            return self.add_content_slide(
                params.get('title', ''),
                params.get('content', [])
            )
        elif slide_type == SlideType.IMAGE:
            return self.add_image_slide(
                params.get('title', ''),
                params.get('image_path', Path('placeholder.jpg')),
                params.get('caption', '')
            )
        elif slide_type == SlideType.SECTION:
            return self.add_section_slide(
                params.get('title', ''),
                params.get('description', '')
            )
        else:
            # Default to content slide
            return self.add_content_slide(
                params.get('title', 'Untitled'),
                params.get('content', [])
            )

    @classmethod
    def from_markdown(cls, markdown_path: Path,
                     template_path: Path | None = None) -> "PresentationBuilder":
        """Create a presentation builder from a markdown template.

        Args:
            markdown_path: Path to markdown template file
            template_path: Optional PowerPoint template file

        Returns:
            Configured presentation builder
        """
        from .markdown import create_presentation_from_markdown

        builder = cls(template_path)
        presentation_template = create_presentation_from_markdown(markdown_path, template_path)
        builder.build_from_template(presentation_template)
        return builder

    def save(self, output_path: Path) -> None:
        """Save the presentation to a file.

        Args:
            output_path: Path where to save the presentation
        """
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(output_path))

    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        return len(self.presentation.slides)

    def get_slide_titles(self) -> list[str]:
        """Get titles of all slides."""
        titles = []
        for slide in self.presentation.slides:
            if slide.shapes.title:
                titles.append(slide.shapes.title.text)
            else:
                titles.append("(No Title)")
        return titles
